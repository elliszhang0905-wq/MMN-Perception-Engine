"""Deterministic vehicle decision snapshots, reports and learning gates.

This module owns the business rules. ``server.py`` only authenticates, assembles
existing cockpit inputs and translates HTTP requests.
"""

from __future__ import annotations

import hashlib
import json
import re
import uuid
from datetime import date, datetime, timezone
from io import BytesIO


SURFACES = (
    ("executive_summary", "高管摘要"),
    ("group_impact", "集团影响"),
    ("sales_warning", "销量预警"),
    ("track_environment", "赛道环境"),
    ("policy_environment", "政策环境"),
    ("communication_momentum", "传播势能"),
    ("platform_position", "平台阵地"),
    ("product_voice", "产品用户之声"),
)
SURFACE_LABELS = dict(SURFACES)
CLAIM_TYPES = {"fact", "inference", "hypothesis", "unknown"}
EVIDENCE_STATUSES = {"aligned", "limited", "conflict", "manual_required", "missing", "stale"}
ACTION_STATUSES = ("proposed", "approved", "planned", "running", "observed", "closed", "cancelled")
LEARNING_STATUSES = {"candidate", "pending_review", "approved", "rejected", "needs_more_observation"}
KNOWHOW_STATUSES = {"pending_review", "approved", "rejected"}


def _now():
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _json(value):
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _load(value, default):
    if value in (None, ""):
        return default
    if isinstance(value, (dict, list)):
        return value
    try:
        return json.loads(value)
    except (TypeError, ValueError, json.JSONDecodeError):
        return default


def _id(prefix):
    return f"{prefix}_{uuid.uuid4().hex}"


def _fingerprint(value):
    return hashlib.sha256(_json(value).encode("utf-8")).hexdigest()


def _required_text(value, label):
    text = str(value or "").strip()
    if not text:
        raise ValueError(f"缺少{label}")
    return text


def _scoped_row(conn, table, item_id, org_id):
    row = conn.execute(f"select * from {table} where id=? and org_id=?", (item_id, org_id)).fetchone()
    if not row:
        raise LookupError("记录不存在或不属于当前组织")
    return row


def init_vehicle_decision_schema(conn):
    conn.executescript(
        """
        create table if not exists vehicle_decision_snapshots (
          id text primary key, org_id text not null, user_id text not null,
          edition text not null, brand text not null default '', model text not null,
          project text not null default '', vehicle_stage text not null default '',
          business_question text not null, core_competitors_json text not null default '[]',
          data_cutoff_at text not null, surface_inputs_json text not null,
          surface_coverage_json text not null, evidence_ids_json text not null default '[]',
          data_fingerprint text not null, source_version text not null,
          created_at text not null
        );
        create index if not exists idx_vehicle_decision_snapshots_scope
          on vehicle_decision_snapshots(org_id, edition, model, created_at desc);

        create table if not exists decision_signals (
          id text primary key, snapshot_id text not null, org_id text not null,
          edition text not null, model text not null, surface text not null,
          time_window_json text not null, claim_type text not null,
          conclusion text not null, evidence_ids_json text not null,
          evidence_status text not null, confidence real, business_impact real,
          recommended_action text not null default '', leading_indicator text not null default '',
          result_indicator text not null default '', stop_condition text not null default '',
          uncertainty text not null default '', source_version text not null,
          generated_at text not null, payload_json text not null default '{}',
          unique(snapshot_id, surface, id)
        );
        create index if not exists idx_decision_signals_snapshot on decision_signals(snapshot_id, surface);

        create table if not exists decision_reports (
          id text primary key, snapshot_id text not null, org_id text not null,
          edition text not null, model text not null, version integer not null,
          status text not null, content_json text not null, content_fingerprint text not null,
          created_by text not null, created_at text not null, published_at text,
          unique(snapshot_id, version)
        );
        create index if not exists idx_decision_reports_scope
          on decision_reports(org_id, edition, model, created_at desc);

        create table if not exists decision_conflicts (
          id text primary key, report_id text not null, snapshot_id text not null,
          org_id text not null, conflict_type text not null, severity text not null,
          signal_ids_json text not null, description text not null,
          status text not null default 'open', resolution text not null default '',
          created_at text not null, resolved_at text
        );

        create table if not exists human_adjudications (
          id text primary key, org_id text not null, target_type text not null,
          target_id text not null, decision text not null, reason text not null,
          adjudicator text not null, created_at text not null
        );

        create table if not exists action_cycles (
          id text primary key, org_id text not null, user_id text not null,
          edition text not null, model text not null, report_id text not null,
          decision_id text not null, source_insight_ids_json text not null,
          evidence_ids_json text not null, knowledge_ids_json text not null,
          hypothesis text not null, owner text not null, target text not null,
          platform text not null, region text not null, audience text not null,
          start_at text not null, end_at text not null, baseline_json text not null,
          target_value_json text not null, leading_indicator text not null,
          result_indicator text not null, stop_condition text not null,
          status text not null, approved_by text not null, approved_at text not null,
          created_at text not null, updated_at text not null,
          unique(org_id, report_id, decision_id)
        );

        create table if not exists action_results (
          id text primary key, action_id text not null, org_id text not null,
          version integer not null, metrics_json text not null,
          actual_execution text not null, completion_rate real,
          actual_time_window text not null, external_variables_json text not null,
          observation text not null, evidence_ids_json text not null,
          recorded_by text not null, created_at text not null, updated_at text not null,
          unique(action_id, version)
        );

        create table if not exists learning_candidates (
          id text primary key, action_id text not null, result_id text not null,
          org_id text not null, edition text not null, model text not null,
          status text not null, hypothesis text not null, actual_action text not null,
          observed_result_json text not null, supported integer,
          counter_evidence text not null, alternative_explanations_json text not null,
          applicability text not null, non_applicability text not null,
          evidence_ids_json text not null, review_reason text not null default '',
          reviewed_by text, reviewed_at text, created_at text not null,
          unique(action_id, result_id)
        );

        create table if not exists knowhow_candidates (
          id text primary key, org_id text not null, edition text not null,
          model text not null, status text not null, learning_ids_json text not null,
          trigger_conditions text not null, applicability text not null,
          non_applicability text not null, execution_steps_json text not null,
          owner text not null, resources_json text not null,
          leading_indicator text not null, result_indicator text not null,
          acceptance_condition text not null, stop_condition text not null,
          ellis_waiver integer not null default 0, waiver_reason text not null default '',
          review_reason text not null default '', reviewed_by text, reviewed_at text,
          created_at text not null
        );

        create table if not exists decision_evidence_links (
          id text primary key, org_id text not null, subject_type text not null,
          subject_id text not null, evidence_id text not null, created_at text not null,
          unique(org_id, subject_type, subject_id, evidence_id)
        );
        """
    )
    conn.commit()


def _signal_payload(snapshot_id, org_id, edition, model, surface, raw, source_version, generated_at):
    raw = dict(raw or {})
    evidence_ids = [str(item).strip() for item in raw.get("evidenceIds") or [] if str(item).strip()]
    claim_type = str(raw.get("claimType") or ("unknown" if not raw else "inference"))
    if claim_type not in CLAIM_TYPES:
        claim_type = "unknown"
    conclusion = str(raw.get("conclusion") or "").strip()
    if not conclusion:
        conclusion = f"{SURFACE_LABELS[surface]}当前没有可用于该车型的决策证据"
        claim_type = "unknown"
    evidence_status = str(raw.get("evidenceStatus") or "")
    if evidence_status not in EVIDENCE_STATUSES:
        if not raw or (claim_type == "unknown" and not evidence_ids):
            evidence_status = "missing"
        elif not evidence_ids:
            evidence_status = "manual_required"
        else:
            evidence_status = "aligned"
    if claim_type == "fact" and not evidence_ids:
        evidence_status = "manual_required"
    stable = {
        "snapshotId": snapshot_id, "orgId": org_id, "edition": edition, "vehicleId": raw.get("vehicleId") or model,
        "model": model, "surface": surface, "timeWindow": raw.get("timeWindow") or "unknown",
        "claimType": claim_type, "conclusion": conclusion, "evidenceIds": evidence_ids,
        "evidenceStatus": evidence_status, "confidence": raw.get("confidence"),
        "businessImpact": raw.get("businessImpact"), "recommendedAction": str(raw.get("recommendedAction") or "补齐证据后由人工选择动作"),
        "leadingIndicator": str(raw.get("leadingIndicator") or ""), "resultIndicator": str(raw.get("resultIndicator") or ""),
        "stopCondition": str(raw.get("stopCondition") or ""), "uncertainty": str(raw.get("uncertainty") or ("证据缺失" if evidence_status == "missing" else "")),
        "sourceVersion": str(raw.get("sourceVersion") or source_version), "generatedAt": generated_at,
        "metricDefinition": str(raw.get("metricDefinition") or ""),
    }
    stable["signalId"] = "sig_" + _fingerprint({"snapshot": snapshot_id, "surface": surface, "payload": stable})[:24]
    return stable


def _snapshot_payload(row, signals):
    return {
        "id": row["id"], "orgId": row["org_id"], "edition": row["edition"], "brand": row["brand"],
        "model": row["model"], "project": row["project"], "vehicleStage": row["vehicle_stage"],
        "businessQuestion": row["business_question"], "coreCompetitors": _load(row["core_competitors_json"], []),
        "dataCutoffAt": row["data_cutoff_at"], "surfaceCoverage": _load(row["surface_coverage_json"], []),
        "evidenceIds": _load(row["evidence_ids_json"], []), "dataFingerprint": row["data_fingerprint"],
        "sourceVersion": row["source_version"], "createdAt": row["created_at"], "signals": signals,
    }


def create_snapshot(conn, body, *, org_id, user_id, edition):
    org_id = _required_text(org_id, "组织")
    model = _required_text(body.get("model"), "车型")
    question = _required_text(body.get("businessQuestion"), "业务问题")
    cutoff = _required_text(body.get("dataCutoffAt"), "报告数据截止时间")
    snapshot_id, created_at = _id("vds"), _now()
    source_version = str(body.get("sourceVersion") or "vehicle-decision-v1")
    inputs = body.get("surfaceInputs") or {}
    signals = []
    for surface, _label in SURFACES:
        items = inputs.get(surface) or []
        if not isinstance(items, list):
            items = [items]
        if not items:
            items = [{}]
        for raw in items:
            signals.append(_signal_payload(snapshot_id, org_id, edition, model, surface, raw, source_version, created_at))
    evidence_ids = sorted({item for signal in signals for item in signal["evidenceIds"]})
    coverage = []
    for surface, label in SURFACES:
        related = [item for item in signals if item["surface"] == surface]
        statuses = {item["evidenceStatus"] for item in related}
        status = "missing" if statuses == {"missing"} else "manual_required" if "manual_required" in statuses else "conflict" if "conflict" in statuses else "limited" if "limited" in statuses else "stale" if "stale" in statuses else "aligned"
        coverage.append({"surface": surface, "label": label, "evidenceStatus": status, "signalCount": len([item for item in related if item["evidenceStatus"] != "missing"]), "timeWindows": sorted({str(item["timeWindow"]) for item in related})})
    frozen = {"metadata": {key: body.get(key) for key in ("brand", "model", "project", "vehicleStage", "businessQuestion", "coreCompetitors", "dataCutoffAt")}, "surfaceInputs": inputs, "signals": signals}
    fingerprint = _fingerprint(frozen)
    conn.execute(
        """insert into vehicle_decision_snapshots
        (id,org_id,user_id,edition,brand,model,project,vehicle_stage,business_question,core_competitors_json,data_cutoff_at,surface_inputs_json,surface_coverage_json,evidence_ids_json,data_fingerprint,source_version,created_at)
        values (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
        (snapshot_id, org_id, user_id, edition, str(body.get("brand") or ""), model, str(body.get("project") or ""), str(body.get("vehicleStage") or ""), question, _json(body.get("coreCompetitors") or []), cutoff, _json(inputs), _json(coverage), _json(evidence_ids), fingerprint, source_version, created_at),
    )
    for signal in signals:
        conn.execute(
            """insert into decision_signals
            (id,snapshot_id,org_id,edition,model,surface,time_window_json,claim_type,conclusion,evidence_ids_json,evidence_status,confidence,business_impact,recommended_action,leading_indicator,result_indicator,stop_condition,uncertainty,source_version,generated_at,payload_json)
            values (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (signal["signalId"], snapshot_id, org_id, edition, model, signal["surface"], _json(signal["timeWindow"]), signal["claimType"], signal["conclusion"], _json(signal["evidenceIds"]), signal["evidenceStatus"], signal["confidence"], signal["businessImpact"], signal["recommendedAction"], signal["leadingIndicator"], signal["resultIndicator"], signal["stopCondition"], signal["uncertainty"], signal["sourceVersion"], signal["generatedAt"], _json(signal)),
        )
        for evidence_id in signal["evidenceIds"]:
            conn.execute("insert or ignore into decision_evidence_links values (?,?,?,?,?,?)", (_id("del"), org_id, "signal", signal["signalId"], evidence_id, created_at))
    conn.commit()
    return get_snapshot(conn, snapshot_id, org_id=org_id)


def get_snapshot(conn, snapshot_id, *, org_id):
    row = _scoped_row(conn, "vehicle_decision_snapshots", snapshot_id, org_id)
    signals = [_load(item["payload_json"], {}) for item in conn.execute("select * from decision_signals where snapshot_id=? and org_id=? order by rowid", (snapshot_id, org_id)).fetchall()]
    return _snapshot_payload(row, signals)


def list_snapshots(conn, *, org_id, edition="china", model="", limit=50):
    limit = max(1, min(int(limit or 50), 200))
    rows = conn.execute("""select * from vehicle_decision_snapshots where org_id=? and edition=? and (?='' or model=?) order by created_at desc limit ?""", (org_id, edition, model, model, limit)).fetchall()
    return [_snapshot_payload(row, []) for row in rows]


def _find_conflicts(signals):
    findings = []
    windows = {str(item["timeWindow"]) for item in signals if str(item["timeWindow"]) not in {"", "unknown"}}
    if len(windows) > 1:
        findings.append({"type": "time_window_incomparable", "severity": "warning", "signalIds": [item["signalId"] for item in signals], "description": "八类表面存在不同数据周期，报告已分别保留，不能直接按同周期比较。"})
    negative = [item for item in signals if re.search(r"风险高|走弱|不足|偏弱|下降|恶化", item["conclusion"])]
    positive = [item for item in signals if re.search(r"正向|增强|领先|改善|增长", item["conclusion"])]
    if negative and positive:
        findings.append({"type": "cross_surface_direction_conflict", "severity": "critical", "signalIds": [negative[0]["signalId"], positive[0]["signalId"]], "description": "不同表面对同一车型给出方向相反的判断，需要人工说明口径或裁决。"})
    for item in signals:
        if item["claimType"] != "fact" and re.search(r"导致|证明|必然|一定会", item["conclusion"]):
            findings.append({"type": "causal_overreach", "severity": "critical", "signalIds": [item["signalId"]], "description": "非事实信号包含因果或必然表述，必须降级或补充因果证据。"})
        if item["claimType"] == "fact" and not item["evidenceIds"]:
            findings.append({"type": "fact_without_evidence", "severity": "critical", "signalIds": [item["signalId"]], "description": "事实信号缺少证据ID。"})
        if item["evidenceStatus"] in {"conflict", "manual_required"}:
            findings.append({"type": "surface_manual_required", "severity": "critical", "signalIds": [item["signalId"]], "description": f"{SURFACE_LABELS[item['surface']]}需要人工复核。"})
    seen = {}
    for item in signals:
        key = re.sub(r"\W+", "", item["conclusion"])
        if key in seen and seen[key] != item["signalId"]:
            findings.append({"type": "duplicate_conclusion", "severity": "info", "signalIds": [seen[key], item["signalId"]], "description": "重复结论已在排序层归并。"})
        seen[key] = item["signalId"]
    return findings


def _report_row(row, conflicts=None):
    return {"id": row["id"], "snapshotId": row["snapshot_id"], "orgId": row["org_id"], "edition": row["edition"], "model": row["model"], "version": row["version"], "status": row["status"], "content": _load(row["content_json"], {}), "contentFingerprint": row["content_fingerprint"], "createdAt": row["created_at"], "publishedAt": row["published_at"], "conflicts": conflicts or []}


def generate_report(conn, snapshot_id, *, org_id, user_id):
    snapshot = get_snapshot(conn, snapshot_id, org_id=org_id)
    signals = snapshot["signals"]
    conflicts = _find_conflicts(signals)
    unresolved_critical = any(item["severity"] == "critical" for item in conflicts)
    ranked = sorted((item for item in signals if item["evidenceStatus"] != "missing"), key=lambda item: (float(item.get("businessImpact") or 0), float(item.get("confidence") or 0)), reverse=True)
    knowledge = list_approved_knowledge(conn, org_id=org_id, edition=snapshot["edition"], model=snapshot["model"])
    content = {
        "executiveConclusion": ranked[0]["conclusion"] if ranked else "当前证据不足，不能形成正式车型结论",
        "scope": {key: snapshot[key] for key in ("brand", "model", "project", "vehicleStage", "businessQuestion", "coreCompetitors", "dataCutoffAt", "dataFingerprint")},
        "surfaceCoverage": snapshot["surfaceCoverage"],
        "topDecisions": ranked[:3],
        "topRisks": [item for item in ranked if re.search(r"风险|走弱|不足|偏弱|下降", item["conclusion"])][:3],
        "topOpportunities": [item for item in ranked if item["claimType"] in {"fact", "inference"} and item["evidenceStatus"] == "aligned"][:3],
        "evidenceChain": [{"signalId": item["signalId"], "surface": item["surface"], "conclusion": item["conclusion"], "claimType": item["claimType"], "evidenceIds": item["evidenceIds"], "timeWindow": item["timeWindow"]} for item in ranked],
        "conflictsAndUnknowns": conflicts + [{"type": "missing_surface", "severity": "warning", "signalIds": [item["signalId"]], "description": item["conclusion"]} for item in signals if item["evidenceStatus"] == "missing"],
        "referencedLearnings": knowledge["learnings"], "referencedKnowhows": knowledge["knowhows"],
        "manualAdjudications": [item for item in conflicts if item["severity"] == "critical"],
        "actions": [], "resultStatus": "尚未建立Action", "learningCandidateStatus": "尚未生成", "knowhowCandidateStatus": "尚未生成",
        "evidenceAppendix": sorted({evidence for item in signals for evidence in item["evidenceIds"]}),
    }
    version = int(conn.execute("select coalesce(max(version),0)+1 from decision_reports where snapshot_id=?", (snapshot_id,)).fetchone()[0])
    report_id, created_at = _id("vdr"), _now()
    status = "manual_required" if unresolved_critical else "draft"
    conn.execute("insert into decision_reports values (?,?,?,?,?,?,?,?,?,?,?,?)", (report_id, snapshot_id, org_id, snapshot["edition"], snapshot["model"], version, status, _json(content), _fingerprint(content), user_id, created_at, None))
    saved_conflicts = []
    for item in conflicts:
        conflict_id = _id("vdc")
        conn.execute("insert into decision_conflicts values (?,?,?,?,?,?,?,?,?,?,?,?)", (conflict_id, report_id, snapshot_id, org_id, item["type"], item["severity"], _json(item["signalIds"]), item["description"], "open", "", created_at, None))
        saved_conflicts.append({"id": conflict_id, **item, "status": "open"})
    conn.commit()
    return _report_row(_scoped_row(conn, "decision_reports", report_id, org_id), saved_conflicts)


def get_report(conn, report_id, *, org_id):
    row = _scoped_row(conn, "decision_reports", report_id, org_id)
    conflicts = [{"id": item["id"], "type": item["conflict_type"], "severity": item["severity"], "signalIds": _load(item["signal_ids_json"], []), "description": item["description"], "status": item["status"], "resolution": item["resolution"]} for item in conn.execute("select * from decision_conflicts where report_id=? and org_id=? order by created_at", (report_id, org_id)).fetchall()]
    return _report_row(row, conflicts)


def list_report_versions(conn, snapshot_id, *, org_id):
    rows = conn.execute("select * from decision_reports where snapshot_id=? and org_id=? order by version desc", (snapshot_id, org_id)).fetchall()
    return [_report_row(row) for row in rows]


def adjudicate_conflict(conn, conflict_id, decision, reason, *, org_id, user_id):
    conflict = _scoped_row(conn, "decision_conflicts", conflict_id, org_id)
    decision = _required_text(decision, "裁决决定")
    reason = _required_text(reason, "裁决理由")
    resolved_at = _now()
    conn.execute("update decision_conflicts set status='resolved', resolution=?, resolved_at=? where id=? and org_id=?", (_json({"decision": decision, "reason": reason, "adjudicator": user_id}), resolved_at, conflict_id, org_id))
    conn.execute("insert into human_adjudications values (?,?,?,?,?,?,?,?)", (_id("adj"), org_id, "decision_conflict", conflict_id, decision, reason, user_id, resolved_at))
    remaining = conn.execute("select count(*) from decision_conflicts where report_id=? and org_id=? and severity='critical' and status!='resolved'", (conflict["report_id"], org_id)).fetchone()[0]
    if not remaining:
        conn.execute("update decision_reports set status='draft' where id=? and org_id=? and status='manual_required'", (conflict["report_id"], org_id))
    conn.commit()
    return get_report(conn, conflict["report_id"], org_id=org_id)


def publish_report(conn, report_id, *, org_id, user_id, approval_note):
    report = get_report(conn, report_id, org_id=org_id)
    if report["status"] == "manual_required" or any(item["severity"] == "critical" and item["status"] != "resolved" for item in report["conflicts"]):
        raise ValueError("关键冲突尚未人工裁决，不能发布报告")
    note = _required_text(approval_note, "人工批准说明")
    published_at = _now()
    conn.execute("update decision_reports set status='published', published_at=? where id=? and org_id=?", (published_at, report_id, org_id))
    conn.execute("insert into human_adjudications values (?,?,?,?,?,?,?,?)", (_id("adj"), org_id, "decision_report", report_id, "published", note, user_id, published_at))
    conn.commit()
    return get_report(conn, report_id, org_id=org_id)


def _action_payload(row):
    return {"id": row["id"], "orgId": row["org_id"], "edition": row["edition"], "model": row["model"], "reportId": row["report_id"], "decisionId": row["decision_id"], "sourceInsightIds": _load(row["source_insight_ids_json"], []), "evidenceIds": _load(row["evidence_ids_json"], []), "knowledgeIds": _load(row["knowledge_ids_json"], []), "hypothesis": row["hypothesis"], "owner": row["owner"], "target": row["target"], "platform": row["platform"], "region": row["region"], "audience": row["audience"], "startAt": row["start_at"], "endAt": row["end_at"], "baseline": _load(row["baseline_json"], {}), "targetValue": _load(row["target_value_json"], {}), "leadingIndicator": row["leading_indicator"], "resultIndicator": row["result_indicator"], "stopCondition": row["stop_condition"], "status": row["status"], "approvedBy": row["approved_by"], "approvedAt": row["approved_at"], "createdAt": row["created_at"], "updatedAt": row["updated_at"]}


def create_action(conn, body, *, org_id, user_id, edition):
    report = get_report(conn, _required_text(body.get("reportId"), "来源报告"), org_id=org_id)
    if report["status"] != "published":
        raise ValueError("报告必须先经人工批准并发布，才能建立Action")
    decision_id = _required_text(body.get("decisionId"), "来源Decision")
    decision = next((item for item in report["content"].get("topDecisions", []) if item.get("signalId") == decision_id), None)
    if not decision:
        raise ValueError("所选Decision不在该报告的可执行决策中")
    if decision.get("evidenceStatus") in {"conflict", "manual_required", "missing"}:
        raise ValueError("该Decision仍需人工裁决，不能直接执行")
    required = {key: _required_text(body.get(key), label) for key, label in (("hypothesis", "行动假设"), ("owner", "负责人"), ("target", "执行对象"), ("platform", "平台"), ("region", "区域"), ("audience", "人群"), ("startAt", "开始时间"), ("endAt", "结束时间"), ("leadingIndicator", "领先指标"), ("resultIndicator", "结果指标"), ("stopCondition", "停止条件"))}
    if edition and edition != report["edition"]:
        raise ValueError("Action版本与来源报告不一致")
    if (_parse_date(required["startAt"]) is None or _parse_date(required["endAt"]) is None
            or _parse_date(required["startAt"]) > _parse_date(required["endAt"])):
        raise ValueError("Action开始和结束时间无效")
    baseline = body.get("baseline")
    if not isinstance(baseline, dict) or not baseline:
        raise ValueError("缺少执行前基线")
    target_value = body.get("targetValue")
    if not isinstance(target_value, dict) or not target_value:
        raise ValueError("缺少目标值")
    action_id, created_at = _id("act"), _now()
    try:
        conn.execute("""insert into action_cycles
            (id,org_id,user_id,edition,model,report_id,decision_id,source_insight_ids_json,
             evidence_ids_json,knowledge_ids_json,hypothesis,owner,target,platform,region,audience,
             start_at,end_at,baseline_json,target_value_json,leading_indicator,result_indicator,
             stop_condition,status,approved_by,approved_at,created_at,updated_at)
            values (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (action_id, org_id, user_id, report["edition"], report["model"], report["id"], decision_id,
             _json([decision_id]), _json(decision.get("evidenceIds") or []),
             _json([item.get("id") for item in report["content"].get("referencedLearnings", []) + report["content"].get("referencedKnowhows", []) if item.get("id")]),
             required["hypothesis"], required["owner"], required["target"], required["platform"],
             required["region"], required["audience"], required["startAt"], required["endAt"],
             _json(baseline), _json(target_value), required["leadingIndicator"], required["resultIndicator"],
             required["stopCondition"], "approved", user_id, created_at, created_at, created_at))
    except Exception as exc:
        if "unique" in str(exc).lower():
            row = conn.execute("select * from action_cycles where org_id=? and report_id=? and decision_id=?", (org_id, report["id"], decision_id)).fetchone()
            return _action_payload(row)
        raise
    conn.commit()
    return _action_payload(_scoped_row(conn, "action_cycles", action_id, org_id))


def update_action_status(conn, action_id, status, *, org_id, user_id):
    row = _scoped_row(conn, "action_cycles", action_id, org_id)
    status = str(status or "")
    if status not in ACTION_STATUSES:
        raise ValueError("Action状态无效")
    current = ACTION_STATUSES.index(row["status"])
    target = ACTION_STATUSES.index(status)
    if status != "cancelled" and target < current:
        raise ValueError("Action状态不能倒退")
    updated_at = _now()
    conn.execute("update action_cycles set status=?, updated_at=? where id=? and org_id=?", (status, updated_at, action_id, org_id))
    conn.execute("insert into human_adjudications values (?,?,?,?,?,?,?,?)", (_id("adj"), org_id, "action", action_id, status, "人工更新Action状态", user_id, updated_at))
    conn.commit()
    return _action_payload(_scoped_row(conn, "action_cycles", action_id, org_id))


def record_result(conn, action_id, body, *, org_id, user_id):
    action = _action_payload(_scoped_row(conn, "action_cycles", action_id, org_id))
    if action["status"] not in {"running", "observed", "closed"}:
        raise ValueError("Action尚未实际执行，不能记录Result")
    metrics = body.get("metrics") or {}
    allowed = ("volume", "interaction", "qualifiedComments", "nsr", "leads", "testDrives", "inquiries", "orders", "sales")
    normalized = {key: metrics[key] if key in metrics else None for key in allowed}
    for key, value in normalized.items():
        if value is not None and not isinstance(value, (int, float)):
            raise ValueError(f"Result指标{key}必须为数字或空值")
    completion_rate = body.get("completionRate")
    if completion_rate is not None and (not isinstance(completion_rate, (int, float)) or not 0 <= completion_rate <= 1):
        raise ValueError("Action完成度必须为0到1之间的数字或空值")
    result_id, created_at = _id("res"), _now()
    version = int(conn.execute("select coalesce(max(version),0)+1 from action_results where action_id=?", (action_id,)).fetchone()[0])
    conn.execute("insert into action_results values (?,?,?,?,?,?,?,?,?,?,?,?,?,?)", (result_id, action_id, org_id, version, _json(normalized), _required_text(body.get("actualExecution"), "实际执行内容"), completion_rate, _required_text(body.get("actualTimeWindow"), "实际时间范围"), _json(body.get("externalVariables") or []), str(body.get("observation") or ""), _json(body.get("evidenceIds") or []), user_id, created_at, created_at))
    conn.commit()
    row = _scoped_row(conn, "action_results", result_id, org_id)
    return {"id": row["id"], "actionId": action_id, "version": version, "metrics": normalized, "actualExecution": row["actual_execution"], "completionRate": row["completion_rate"], "actualTimeWindow": row["actual_time_window"], "externalVariables": _load(row["external_variables_json"], []), "observation": row["observation"], "evidenceIds": _load(row["evidence_ids_json"], []), "recordedBy": row["recorded_by"], "createdAt": row["created_at"]}


def _parse_date(value):
    try:
        return date.fromisoformat(str(value or "")[:10])
    except ValueError:
        return None


def generate_learning_candidate(conn, action_id, body, *, org_id, user_id, as_of=None):
    action = _action_payload(_scoped_row(conn, "action_cycles", action_id, org_id))
    if action["status"] not in {"observed", "closed"}:
        raise ValueError("Action尚未完成观察，不能生成Learning Candidate")
    if not action["baseline"]:
        raise ValueError("缺少执行前基线")
    result = conn.execute("select * from action_results where action_id=? and org_id=? order by version desc limit 1", (action_id, org_id)).fetchone()
    if not result:
        raise ValueError("没有真实Result，不能生成Learning Candidate")
    end_at, observed_at = _parse_date(action["endAt"]), _parse_date(as_of or _now())
    status = "needs_more_observation" if end_at and observed_at and observed_at < end_at else "pending_review"
    counter = _required_text(body.get("counterEvidence"), "反证")
    alternatives = body.get("alternativeExplanations") or []
    if not alternatives:
        raise ValueError("至少记录一种其他可能解释")
    candidate_id, created_at = _id("lrn"), _now()
    evidence_ids = sorted(set(action["evidenceIds"] + _load(result["evidence_ids_json"], [])))
    try:
        conn.execute("""insert into learning_candidates
            (id,action_id,result_id,org_id,edition,model,status,hypothesis,actual_action,
             observed_result_json,supported,counter_evidence,alternative_explanations_json,
             applicability,non_applicability,evidence_ids_json,review_reason,reviewed_by,reviewed_at,created_at)
            values (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (candidate_id, action_id, result["id"], org_id, action["edition"], action["model"], status,
             action["hypothesis"], result["actual_execution"], result["metrics_json"],
             None if body.get("supported") is None else int(bool(body.get("supported"))), counter,
             _json(alternatives), _required_text(body.get("applicability"), "适用条件"),
             _required_text(body.get("nonApplicability"), "不适用条件"), _json(evidence_ids),
             "", None, None, created_at))
    except Exception as exc:
        if "unique" in str(exc).lower():
            row = conn.execute("select * from learning_candidates where action_id=? and result_id=?", (action_id, result["id"])).fetchone()
            return _learning_payload(row)
        raise
    conn.commit()
    return _learning_payload(_scoped_row(conn, "learning_candidates", candidate_id, org_id))


def _learning_payload(row):
    return {"id": row["id"], "actionId": row["action_id"], "resultId": row["result_id"], "orgId": row["org_id"], "edition": row["edition"], "model": row["model"], "status": row["status"], "hypothesis": row["hypothesis"], "actualAction": row["actual_action"], "observedResult": _load(row["observed_result_json"], {}), "supported": None if row["supported"] is None else bool(row["supported"]), "counterEvidence": row["counter_evidence"], "alternativeExplanations": _load(row["alternative_explanations_json"], []), "applicability": row["applicability"], "nonApplicability": row["non_applicability"], "evidenceIds": _load(row["evidence_ids_json"], []), "reviewReason": row["review_reason"], "reviewedBy": row["reviewed_by"], "reviewedAt": row["reviewed_at"], "createdAt": row["created_at"]}


def review_learning_candidate(conn, candidate_id, decision, reason, *, org_id, user_id):
    row = _scoped_row(conn, "learning_candidates", candidate_id, org_id)
    decision = str(decision or "")
    if decision not in {"approved", "rejected", "needs_more_observation"}:
        raise ValueError("Learning人工裁决状态无效")
    if row["status"] == "needs_more_observation" and decision == "approved":
        raise ValueError("观察周期尚未结束，不能批准Learning")
    reviewed_at = _now()
    supported = 1 if decision == "approved" else 0 if decision == "rejected" else None
    conn.execute("update learning_candidates set status=?,supported=?,review_reason=?,reviewed_by=?,reviewed_at=? where id=? and org_id=?", (decision, supported, _required_text(reason, "裁决理由"), user_id, reviewed_at, candidate_id, org_id))
    conn.execute("insert into human_adjudications values (?,?,?,?,?,?,?,?)", (_id("adj"), org_id, "learning_candidate", candidate_id, decision, reason, user_id, reviewed_at))
    conn.commit()
    return _learning_payload(_scoped_row(conn, "learning_candidates", candidate_id, org_id))


def _knowhow_payload(row):
    return {"id": row["id"], "orgId": row["org_id"], "edition": row["edition"], "model": row["model"], "status": row["status"], "learningIds": _load(row["learning_ids_json"], []), "triggerConditions": row["trigger_conditions"], "applicability": row["applicability"], "nonApplicability": row["non_applicability"], "executionSteps": _load(row["execution_steps_json"], []), "owner": row["owner"], "resources": _load(row["resources_json"], []), "leadingIndicator": row["leading_indicator"], "resultIndicator": row["result_indicator"], "acceptanceCondition": row["acceptance_condition"], "stopCondition": row["stop_condition"], "ellisWaiver": bool(row["ellis_waiver"]), "waiverReason": row["waiver_reason"], "reviewReason": row["review_reason"], "reviewedBy": row["reviewed_by"], "reviewedAt": row["reviewed_at"], "createdAt": row["created_at"]}


def generate_knowhow_candidate(conn, body, *, org_id, user_id):
    learning_ids = list(dict.fromkeys(str(item) for item in body.get("learningIds") or [] if str(item)))
    if not learning_ids:
        raise ValueError("至少选择一个已批准Learning")
    placeholders = ",".join("?" for _ in learning_ids)
    rows = conn.execute(f"select * from learning_candidates where id in ({placeholders}) and org_id=? and status='approved'", (*learning_ids, org_id)).fetchall()
    if len(rows) != len(learning_ids):
        raise ValueError("Know-how只能来源于当前组织已批准的Learning")
    waiver = body.get("ellisWaiver") is True
    if len(rows) < 2 and not waiver:
        raise ValueError("Know-how默认需要至少两个窗口、区域或车型的重复验证")
    waiver_reason = _required_text(body.get("waiverReason"), "Ellis人工豁免理由") if waiver else ""
    first = _learning_payload(rows[0])
    candidate_id, created_at = _id("khw"), _now()
    values = {
        "triggerConditions": str(body.get("triggerConditions") or first["hypothesis"]),
        "applicability": str(body.get("applicability") or first["applicability"]),
        "nonApplicability": str(body.get("nonApplicability") or first["nonApplicability"]),
        "executionSteps": body.get("executionSteps") or [first["actualAction"]],
        "owner": str(body.get("owner") or "项目负责人"), "resources": body.get("resources") or ["来源Action与Result证据"],
        "leadingIndicator": str(body.get("leadingIndicator") or "沿用来源Action领先指标"),
        "resultIndicator": str(body.get("resultIndicator") or "沿用来源Action结果指标"),
        "acceptanceCondition": str(body.get("acceptanceCondition") or "重复验证结果支持原假设"),
        "stopCondition": str(body.get("stopCondition") or "出现边界外反例或结果指标恶化"),
    }
    conn.execute("""insert into knowhow_candidates
        (id,org_id,edition,model,status,learning_ids_json,trigger_conditions,applicability,
         non_applicability,execution_steps_json,owner,resources_json,leading_indicator,
         result_indicator,acceptance_condition,stop_condition,ellis_waiver,waiver_reason,
         review_reason,reviewed_by,reviewed_at,created_at)
        values (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
        (candidate_id, org_id, first["edition"], first["model"], "pending_review", _json(learning_ids),
         values["triggerConditions"], values["applicability"], values["nonApplicability"],
         _json(values["executionSteps"]), values["owner"], _json(values["resources"]),
         values["leadingIndicator"], values["resultIndicator"], values["acceptanceCondition"],
         values["stopCondition"], int(waiver), waiver_reason, "", None, None, created_at))
    conn.commit()
    return _knowhow_payload(_scoped_row(conn, "knowhow_candidates", candidate_id, org_id))


def review_knowhow_candidate(conn, candidate_id, decision, reason, *, org_id, user_id):
    _scoped_row(conn, "knowhow_candidates", candidate_id, org_id)
    if decision not in {"approved", "rejected"}:
        raise ValueError("Know-how人工裁决状态无效")
    reviewed_at = _now()
    conn.execute("update knowhow_candidates set status=?,review_reason=?,reviewed_by=?,reviewed_at=? where id=? and org_id=?", (decision, _required_text(reason, "裁决理由"), user_id, reviewed_at, candidate_id, org_id))
    conn.execute("insert into human_adjudications values (?,?,?,?,?,?,?,?)", (_id("adj"), org_id, "knowhow_candidate", candidate_id, decision, reason, user_id, reviewed_at))
    conn.commit()
    return _knowhow_payload(_scoped_row(conn, "knowhow_candidates", candidate_id, org_id))


def list_approved_knowledge(conn, *, org_id, edition, model):
    learnings = [_learning_payload(row) for row in conn.execute("select * from learning_candidates where org_id=? and edition=? and model=? and status='approved' order by reviewed_at desc", (org_id, edition, model)).fetchall()]
    knowhows = [_knowhow_payload(row) for row in conn.execute("select * from knowhow_candidates where org_id=? and edition=? and model=? and status='approved' order by reviewed_at desc", (org_id, edition, model)).fetchall()]
    return {"learnings": learnings, "knowhows": knowhows}


def get_flow_state(conn, *, org_id, edition, model):
    actions = [_action_payload(row) for row in conn.execute(
        "select * from action_cycles where org_id=? and edition=? and model=? order by updated_at desc",
        (org_id, edition, model),
    ).fetchall()]
    results = []
    if actions:
        action_ids = [item["id"] for item in actions]
        placeholders = ",".join("?" for _ in action_ids)
        for row in conn.execute(
            f"select * from action_results where org_id=? and action_id in ({placeholders}) order by created_at desc",
            (org_id, *action_ids),
        ).fetchall():
            results.append({
                "id": row["id"], "actionId": row["action_id"], "version": row["version"],
                "metrics": _load(row["metrics_json"], {}), "actualExecution": row["actual_execution"],
                "completionRate": row["completion_rate"], "actualTimeWindow": row["actual_time_window"],
                "externalVariables": _load(row["external_variables_json"], []), "observation": row["observation"],
                "evidenceIds": _load(row["evidence_ids_json"], []), "recordedBy": row["recorded_by"],
                "createdAt": row["created_at"], "updatedAt": row["updated_at"],
            })
    learnings = [_learning_payload(row) for row in conn.execute(
        "select * from learning_candidates where org_id=? and edition=? and model=? order by created_at desc",
        (org_id, edition, model),
    ).fetchall()]
    knowhows = [_knowhow_payload(row) for row in conn.execute(
        "select * from knowhow_candidates where org_id=? and edition=? and model=? order by created_at desc",
        (org_id, edition, model),
    ).fetchall()]
    return {"actions": actions, "results": results, "learningCandidates": learnings, "knowhowCandidates": knowhows}


def render_report_markdown(report):
    content = report["content"]
    scope = content.get("scope") or {}
    lines = [f"# {scope.get('model') or report.get('model')} 车型综合决策报告", "", f"- 报告版本：{report.get('version')}", f"- 数据截止：{scope.get('dataCutoffAt') or '未知'}", f"- 业务问题：{scope.get('businessQuestion') or '未填写'}", f"- 发布状态：{report.get('status')}", "", "## 高管结论", "", content.get("executiveConclusion") or "当前证据不足。", "", "## 三项重要判断", ""]
    for item in content.get("topDecisions") or []:
        lines.append(f"- [{item.get('claimType')}] {item.get('conclusion')}（{item.get('surface')}，{item.get('timeWindow')}）")
    lines += ["", "## 八类看板覆盖", ""]
    for item in content.get("surfaceCoverage") or []:
        lines.append(f"- {item.get('label')}：{item.get('evidenceStatus')}，有效信号 {item.get('signalCount')} 条")
    lines += ["", "## 冲突、未知与人工裁决", ""]
    for item in content.get("conflictsAndUnknowns") or []:
        lines.append(f"- {item.get('description')}")
    lines += ["", "## 证据与缺口", ""]
    for item in content.get("evidenceChain") or []:
        evidence = ", ".join(item.get("evidenceIds") or []) or "缺失"
        lines.append(f"- {item.get('conclusion')} | 证据：{evidence}")
    return "\n".join(lines) + "\n"


def render_report_pptx(report):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    content, scope = report["content"], report["content"].get("scope") or {}
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    navy, teal, muted = RGBColor(15, 37, 56), RGBColor(20, 132, 126), RGBColor(79, 101, 116)

    def slide(title, subtitle, bullets):
        page = prs.slides.add_slide(prs.slide_layouts[6])
        page.background.fill.solid(); page.background.fill.fore_color.rgb = RGBColor(247, 250, 250)
        head = page.shapes.add_textbox(Inches(.7), Inches(.55), Inches(11.9), Inches(.7)).text_frame
        head.paragraphs[0].text = title; head.paragraphs[0].font.size = Pt(28); head.paragraphs[0].font.bold = True; head.paragraphs[0].font.color.rgb = navy
        sub = page.shapes.add_textbox(Inches(.72), Inches(1.25), Inches(11.8), Inches(.45)).text_frame
        sub.paragraphs[0].text = subtitle; sub.paragraphs[0].font.size = Pt(11); sub.paragraphs[0].font.color.rgb = teal
        body = page.shapes.add_textbox(Inches(.82), Inches(1.95), Inches(11.6), Inches(4.9)).text_frame
        body.word_wrap = True
        for index, text in enumerate(bullets or ["当前没有可发布内容"]):
            p = body.paragraphs[0] if index == 0 else body.add_paragraph(); p.text = str(text)[:420]; p.font.size = Pt(15); p.font.color.rgb = muted; p.space_after = Pt(10)

    slide(f"{scope.get('model') or report.get('model')} 车型综合决策报告", f"版本 {report.get('version')} | 数据截止 {scope.get('dataCutoffAt') or '未知'}", [content.get("executiveConclusion")])
    slide("八类驾驶舱覆盖", "缺失、过期和冲突均显式保留", [f"{item.get('label')}：{item.get('evidenceStatus')}，{item.get('signalCount')} 条有效信号" for item in content.get("surfaceCoverage") or []])
    slide("Top Decisions", "事实、推断、假设和未知分层呈现", [f"[{item.get('claimType')}] {item.get('conclusion')} | {item.get('timeWindow')}" for item in content.get("topDecisions") or []])
    slide("冲突与证据缺口", "关键冲突必须经人工裁决", [item.get("description") for item in content.get("conflictsAndUnknowns") or []])
    slide("Action与结果闭环", "Evidence → Insight → Decision → Action → Result → Learning → Know-how", ["Action只能来自已发布报告中的人工选择Decision", "Result缺失保持为空，不自动填0", "Learning与Know-how均需人工批准后进入下一轮检索"])
    output = BytesIO(); prs.save(output); return output.getvalue()
