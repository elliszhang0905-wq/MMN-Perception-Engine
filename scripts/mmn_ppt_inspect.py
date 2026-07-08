#!/usr/bin/env python3
"""Inspect PPTX text and slide inventory with python-pptx."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", help="PPTX file to inspect.")
    parser.add_argument("--out", default=None, help="Optional JSON output path.")
    return parser.parse_args()


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()
    prs = Presentation(str(pptx_path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        items = []
        for shape in slide.shapes:
            text = shape_text(shape)
            if text:
                items.append(
                    {
                        "name": shape.name,
                        "text": text,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    }
                )
        slides.append({"pageNumber": index, "shapeCount": len(slide.shapes), "textItems": items})

    payload = {"pptx": str(pptx_path), "slideCount": len(prs.slides), "slides": slides}
    rendered = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.out:
        Path(args.out).write_text(rendered, encoding="utf-8")
    print(rendered)


if __name__ == "__main__":
    main()
